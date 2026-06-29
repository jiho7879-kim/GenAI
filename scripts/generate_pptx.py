"""
Generate the Final Project PPTX presentation.
All elements are editable shapes/text boxes — NO embedded images.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──

PRIMARY   = RGBColor(0x1A, 0x1A, 0x2E)  # Dark navy
ACCENT    = RGBColor(0x00, 0x6D, 0x77)  # Teal
ACCENT2   = RGBColor(0x2E, 0x86, 0xAB)  # Light blue
ACCENT3   = RGBColor(0xE9, 0xC4, 0x6A)  # Gold
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xF0, 0xF0, 0xF0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY  = RGBColor(0x66, 0x66, 0x66)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
RED       = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE    = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helper Functions ──

def add_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a rectangle shape (editable)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()  # no line by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a rounded rectangle shape (editable)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_arrow_right(slide, left, top, width, height, fill_color=ACCENT):
    """Add a right-pointing arrow (editable)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with editable text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, line_spacing=1.3, font_name="Calibri"):
    """Add a text box with multiple lines (editable)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, fs, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs if fs else font_size)
        p.font.bold = bold
        p.font.color.rgb = clr if clr else color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox


def add_circle(slide, left, top, size, fill_color=ACCENT):
    """Add a circle/oval shape (editable)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_chevron(slide, left, top, width, height, fill_color=ACCENT):
    """Add a chevron shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def set_shape_text(shape, text, font_size=14, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER, font_name="Calibri"):
    """Set text on a shape with full editability."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    return shape


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, PRIMARY)

# Decorative top bar
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), fill_color=ACCENT3)

# Title
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "Research Agent", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
            "TCAD + Semiconductor Device Paper Assistant", font_size=28, color=ACCENT3, alignment=PP_ALIGN.CENTER)

# Divider line
add_rect(slide, Inches(4.5), Inches(3.8), Inches(4), Inches(0.04), fill_color=ACCENT3)

# Metadata
add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(1.5),
            "Fast Campus Gen AI Intensive Course — Final Project\n"
            "석사 2년용 개인 Research Assistant\n"
            "반도체소자 + TCAD Simulation + ML Methodology",
            font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
            "2026. 06. 29 — 07. 10  |  Ollama + Langchain + FAISS + BGE-m3 + Streamlit",
            font_size=14, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Header bar
add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
            "Problem Statement", font_size=32, bold=True, color=WHITE)

# Left column: Problem
add_rounded_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), fill_color=LIGHT_GRAY)
add_textbox(slide, Inches(1.2), Inches(1.8), Inches(4.8), Inches(0.5),
            "❌ The Problem", font_size=22, bold=True, color=RED)

problems = [
    "석사 2년 동안 수백 편의 논문을 읽고 관리해야 함",
    "반도체소자 + TCAD + ML — 3개 도메인을 동시에 커버",
    "PDF 파일이 흩어져 있고, 원하는 내용을 찾기 어려움",
    "새로운 논문을 찾고 요약하는 데 시간 소모",
    "시험기간(중간 7/2, 기말 7/10)과 프로젝트 병행",
]
lines = [(f"• {p}", False, 15, DARK_GRAY) for p in problems]
add_multiline_textbox(slide, Inches(1.2), Inches(2.5), Inches(4.8), Inches(4.0), lines)

# Right column: Solution
add_rounded_rect(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(5.2), fill_color=RGBColor(0xE8, 0xF8, 0xF5))
add_textbox(slide, Inches(7.4), Inches(1.8), Inches(4.8), Inches(0.5),
            "✅ Our Solution", font_size=22, bold=True, color=GREEN)

solutions = [
    "하나의 Streamlit 앱으로 모든 논문 관리",
    "PDF 업로드 → 자동 인덱싱 → 질의응답",
    "Multi-Agent가 의도를 파악하여 최적의 답변",
    "Arxiv 검색으로 최신 논문 발견",
    "시행착오를 체계적으로 기록 (final report 활용)",
]
lines = [(f"• {s}", False, 15, DARK_GRAY) for s in solutions]
add_multiline_textbox(slide, Inches(7.4), Inches(2.5), Inches(4.8), Inches(4.0), lines)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — DESIGN CONSTRAINTS
# ═══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
            "Design Constraints", font_size=32, bold=True, color=WHITE)

# 4 constraint cards
constraints = [
    ("💻", "CPU-only", "Intel Iris Xe iGPU\n32GB RAM", ACCENT),
    ("🔒", "사외망", "외부 API-key 서비스\n사용 불가", ACCENT2),
    ("💰", "Free", "Ollama Local LLM\n+ Zen Free fallback", GREEN),
    ("🚀", "Single Command", "pip install &&\nstreamlit run app.py", ORANGE),
]

for i, (emoji, title, desc, color) in enumerate(constraints):
    x = Inches(0.8 + i * 3.1)
    y = Inches(1.6)
    card = add_rounded_rect(slide, x, y, Inches(2.8), Inches(3.5), fill_color=LIGHT_GRAY)
    # Top color bar
    add_rect(slide, x, y, Inches(2.8), Inches(0.08), fill_color=color)
    add_textbox(slide, x, y + Inches(0.3), Inches(2.8), Inches(0.6),
                title, font_size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    # Description
    lines = desc.split("\n")
    lines_formatted = [(l, False, 14, DARK_GRAY) for l in lines]
    add_multiline_textbox(slide, x + Inches(0.2), y + Inches(1.0), Inches(2.4), Inches(2.2),
                          lines_formatted, alignment=PP_ALIGN.CENTER)

# Bottom: Impact statement
add_rounded_rect(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2), fill_color=PRIMARY)
add_textbox(slide, Inches(1.2), Inches(5.7), Inches(10.8), Inches(0.8),
            "모든 기술 선택은 위 4가지 제약조건의 교집합에서 이루어짐",
            font_size=18, color=ACCENT3, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — TECH STACK
# ═══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
            "Tech Stack", font_size=32, bold=True, color=WHITE)

# 6 tech cards in 2 rows
techs = [
    ("Streamlit", "Frontend UI", ACCENT),
    ("Langchain", "RAG Framework", ACCENT2),
    ("FAISS (CPU)", "Vector Store", GREEN),
    ("BAAI/bge-m3", "Embedding Model", ORANGE),
    ("Ollama", "Local LLM", RED),
    ("Arxiv API", "Paper Search", RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (name, role, color) in enumerate(techs):
    row = i // 3
    col = i % 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(1.6 + row * 2.6)
    card = add_rounded_rect(slide, x, y, Inches(3.6), Inches(2.2), fill_color=LIGHT_GRAY, line_color=color, line_width=2)
    add_textbox(slide, x, y + Inches(0.3), Inches(3.6), Inches(0.6),
                name, font_size=20, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(1.0), Inches(3.6), Inches(0.6),
                role, font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — SYSTEM ARCHITECTURE (ALL SHAPES, NO IMAGES)
# ═══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
            "System Architecture", font_size=32, bold=True, color=WHITE)

# Flow: User → Supervisor → 3 Agents → Answer

# User box
user_box = add_rounded_rect(slide, Inches(0.8), Inches(3.0), Inches(1.6), Inches(0.9), fill_color=ACCENT)
set_shape_text(user_box, "User\nInput", font_size=14, bold=True, color=WHITE)

# Arrow to Supervisor
add_arrow_right(slide, Inches(2.5), Inches(3.2), Inches(0.8), Inches(0.5), fill_color=MED_GRAY)

# Supervisor box (bigger)
sup_box = add_rounded_rect(slide, Inches(3.5), Inches(2.3), Inches(2.2), Inches(2.2), fill_color=PRIMARY)
set_shape_text(sup_box, "Supervisor\nAgent", font_size=16, bold=True, color=ACCENT3)
tf = sup_box.text_frame
p = tf.add_paragraph()
p.text = "Intent: {paper,\nsearch, report}"
p.font.size = Pt(11)
p.font.color.rgb = MED_GRAY
p.font.name = "Calibri"

# 3 arrows from Supervisor to Agents
add_arrow_right(slide, Inches(5.8), Inches(2.5), Inches(0.6), Inches(0.35), fill_color=ACCENT)
add_arrow_right(slide, Inches(5.8), Inches(3.2), Inches(0.6), Inches(0.35), fill_color=ACCENT2)
add_arrow_right(slide, Inches(5.8), Inches(3.9), Inches(0.6), Inches(0.35), fill_color=GREEN)

# 3 Agent boxes
agents_data = [
    ("Paper\nAgent", "RAG Q&A\nFAISS + BGE-m3\n+ Ollama", Inches(6.6), Inches(2.0), ACCENT),
    ("Arxiv\nSearch\nAgent", "Arxiv API\n+ LLM Summary", Inches(6.6), Inches(2.8), ACCENT2),
    ("Report\nAgent", "Markdown\nFormatting", Inches(6.6), Inches(3.6), GREEN),
]

for name, desc, x, y, color in agents_data:
    box = add_rounded_rect(slide, x, y, Inches(2.0), Inches(1.6), fill_color=LIGHT_GRAY, line_color=color, line_width=2)
    set_shape_text(box, name, font_size=13, bold=True, color=color)
    tf = box.text_frame
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Calibri"

# 3 arrows merging to Answer
add_arrow_right(slide, Inches(8.8), Inches(2.2), Inches(0.6), Inches(0.3), fill_color=MED_GRAY)
add_arrow_right(slide, Inches(8.8), Inches(3.4), Inches(0.6), Inches(0.3), fill_color=MED_GRAY)
add_arrow_right(slide, Inches(8.8), Inches(4.2), Inches(0.6), Inches(0.3), fill_color=MED_GRAY)

# Answer box
ans_box = add_rounded_rect(slide, Inches(9.6), Inches(2.8), Inches(1.8), Inches(1.2), fill_color=ACCENT3)
set_shape_text(ans_box, "Final\nAnswer + Sources", font_size=14, bold=True, color=DARK_GRAY)

# Streamlit at bottom
str_box = add_rounded_rect(slide, Inches(3.5), Inches(5.5), Inches(5.5), Inches(1.0), fill_color=PRIMARY)
set_shape_text(str_box, "Streamlit UI — 3 Tabs: Paper Lab | Arxiv Search | About", font_size=16, bold=True, color=WHITE)

# Arrow from Answer to Streamlit
add_arrow_right(slide, Inches(9.5), Inches(5.0), Inches(1.5), Inches(0.3), fill_color=MED_GRAY)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — KEY FEATURES
# ═══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
            "Key Features", font_size=32, bold=True, color=WHITE)

features = [
    ("📤", "PDF Ingestion", "PyMuPDF → Chunk → BGE-m3 → FAISS\n자동 임베딩 및 인덱싱", ACCENT),
    ("💬", "RAG Q&A", "Langchain RetrievalQA\n+ Ollama LLM + Source Citation", ACCENT2),
    ("🤖", "Multi-Agent", "Supervisor → Paper / Arxiv / Report\n의도 기반 라우팅", GREEN),
    ("🔍", "Arxiv Search", "최신 논문 검색\n+ LLM 요약 생성", ORANGE),
    ("📋", "Report Gen", "RAG + Search 통합 보고서\nMarkdown 자동 포맷팅", RED),
    ("🧠", "Agent Trace", "실시간 Agent 의사결정\n과정 시각화", RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (emoji, title, desc, color) in enumerate(features):
    row = i // 3
    col = i % 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(1.6 + row * 2.8)
    card = add_rounded_rect(slide, x, y, Inches(3.6), Inches(2.4), fill_color=LIGHT_GRAY, line_color=color, line_width=2)
    add_textbox(slide, x, y + Inches(0.2), Inches(3.6), Inches(0.5),
                f"{emoji}  {title}", font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    desc_lines = desc.split("\n")
    formatted = [(l, False, 13, DARK_GRAY) for l in desc_lines]
    add_multiline_textbox(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.2), Inches(1.4),
                          formatted, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 7 — AGENT PIPELINE DETAIL
# ═══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
            "Agent Pipeline Detail", font_size=32, bold=True, color=WHITE)

# Supervisor detail card
add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.8), fill_color=LIGHT_GRAY, line_color=PRIMARY, line_width=1)
add_textbox(slide, Inches(1.2), Inches(1.6), Inches(3), Inches(0.4),
            "Supervisor Agent", font_size=20, bold=True, color=PRIMARY)

sup_lines = [
    ("Ollama qwen3.5:4b, temperature=0.0, num_predict=64", False, 13, MED_GRAY),
    ("Intent → paper: RAG 질의  |  search: Arxiv 검색  |  report: 통합 보고서", False, 13, DARK_GRAY),
    ("Classification prompt with few-shot examples → deterministic label output", False, 13, MED_GRAY),
]
add_multiline_textbox(slide, Inches(1.2), Inches(2.1), Inches(10.5), Inches(1.0), sup_lines)

# 3 agent columns
agent_details = [
    ("Paper Agent", ACCENT,
     ["RAGChain wrapping FAISS + Ollama",
      "RetrievalQA (chain_type='stuff')",
      "Source dedup by content hash",
      "Model switching at runtime",
      "Summarization support"]),
    ("Arxiv Search Agent", ACCENT2,
     ["Arxiv API (no API key needed)",
      "1-hour in-memory cache",
      "LLM-generated paper summaries",
      "Detail fetch by Arxiv ID",
      "Max results configurable"]),
    ("Report Agent", GREEN,
     ["Markdown formatting",
      "RAG + Search fusion reports",
      "Source citation formatting",
      "Optional LLM enhancement",
      "Multi-section report layout"]),
]

for i, (name, color, items) in enumerate(agent_details):
    x = Inches(0.8 + i * 4.0)
    y = Inches(3.6)
    add_rounded_rect(slide, x, y, Inches(3.6), Inches(3.2), fill_color=LIGHT_GRAY, line_color=color, line_width=2)
    add_textbox(slide, x, y + Inches(0.1), Inches(3.6), Inches(0.5),
                name, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    item_lines = [(f"• {item}", False, 13, DARK_GRAY) for item in items]
    add_multiline_textbox(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.2), Inches(2.5), item_lines)


# ═══════════════════════════════════════════════════════════
# SLIDE 8 — DESIGN DECISIONS
# ═══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
            "Key Design Decisions", font_size=32, bold=True, color=WHITE)

decisions = [
    ("Q1", "Vector Store", "FAISS vs Chroma", "FAISS", "서버리스, 단일 명령어 실행"),
    ("Q2", "Embedding", "Ollama vs BGE-m3", "BGE-m3", "다국어 + TCAD 용어, CPU 동작"),
    ("Q3", "RAG Framework", "Langchain vs 직접 구현", "Langchain", "RetrievalQA stuff 체인 효율적"),
    ("Q4", "Agent Pattern", "Langgraph vs Class", "Class-based", "4개 Agent는 단순 composition"),
    ("Q5", "Intent Classify", "LLM vs Rule-based", "LLM", "확장성, 프로젝트 평가 반영"),
    ("Q6", "PDF Parser", "PyMuPDF vs Unstructured", "PyMuPDF", "가벼움, Langchain 호환"),
]

# Table header
header_y = Inches(1.5)
col_widths = [Inches(0.6), Inches(1.6), Inches(2.4), Inches(1.4), Inches(5.0)]
col_starts = [Inches(0.8)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

headers = ["Q", "Category", "Options", "Choice", "Rationale"]
for j, (cs, cw, h) in enumerate(zip(col_starts, col_widths, headers)):
    add_rect(slide, cs, header_y, cw, Inches(0.5), fill_color=PRIMARY)
    add_textbox(slide, cs, header_y + Inches(0.05), cw, Inches(0.4),
                h, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Table rows
for i, (q, cat, opts, choice, rationale) in enumerate(decisions):
    row_y = header_y + Inches(0.5) + Inches(0.45 * i)
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
    for j, (cs, cw, val) in enumerate(zip(col_starts, col_widths, [q, cat, opts, choice, rationale])):
        add_rect(slide, cs, row_y, cw, Inches(0.45), fill_color=bg_color)
        is_bold = (j == 3)  # Bold the "Choice" column
        clr = ACCENT if is_bold else DARK_GRAY
        add_textbox(slide, cs + Inches(0.05), row_y + Inches(0.05), cw - Inches(0.1), Inches(0.35),
                    val, font_size=11, bold=is_bold, color=clr, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 9 — TRIAL & ERROR
# ═══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
            "Trial & Error Documentation", font_size=32, bold=True, color=WHITE)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
            "모든 시행착오는 .omo/errors/에 체계적으로 기록 — Final Report 포함",
            font_size=14, color=MED_GRAY)

errors_data = [
    ("01", "PDF 텍스트 추출 실패", "code", "이미지 PDF는 OCR 필요, PyMuPDF만으로 불가"),
    ("02", "대용량 PDF look_at 타임아웃", "code", "150p PDF → 페이지 단위 분할 처리"),
    ("03", "PowerShell vs Bash 문법", "code", "Windows `&&` → `; if ($?) { }`"),
    ("04", "계획 v1→v2→v3 진화", "decision", "Q1~Q5 기술 선택 근거 문서화"),
    ("05", "아키텍쳐 결정", "decision", "FAISS, BGE-m3, Class-based Agent 등"),
    ("06", "pip install 타임아웃", "code", "torch 2GB+ → timeout 10분으로 증가"),
    ("07", "Langchain 모듈 경로 변경", "code", "chains → langchain_classic migration"),
    ("08", "통합 테스트 결과", "test", "Import 체인 검증 완료"),
]

# Table header
header_y = Inches(1.8)
col_widths2 = [Inches(0.5), Inches(2.8), Inches(0.8), Inches(6.5)]
col_starts2 = [Inches(0.8)]
for w in col_widths2[:-1]:
    col_starts2.append(col_starts2[-1] + w)

headers2 = ["#", "Issue", "Type", "Key Lesson"]
for j, (cs, cw, h) in enumerate(zip(col_starts2, col_widths2, headers2)):
    add_rect(slide, cs, header_y, cw, Inches(0.4), fill_color=PRIMARY)
    add_textbox(slide, cs, header_y + Inches(0.03), cw, Inches(0.35),
                h, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

for i, (num, issue, itype, lesson) in enumerate(errors_data):
    row_y = header_y + Inches(0.4) + Inches(0.5 * i)
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
    for j, (cs, cw, val) in enumerate(zip(col_starts2, col_widths2, [num, issue, itype, lesson])):
        add_rect(slide, cs, row_y, cw, Inches(0.5), fill_color=bg_color)
        color_type = RED if itype == "code" else (GREEN if itype == "decision" else ORANGE)
        clr = color_type if j == 2 else DARK_GRAY
        add_textbox(slide, cs + Inches(0.05), row_y + Inches(0.05), cw - Inches(0.1), Inches(0.4),
                    val, font_size=11, bold=(j == 2), color=clr, alignment=PP_ALIGN.CENTER if j in [0, 2] else PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════
# SLIDE 10 — RUNNING THE APP
# ═══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
            "How to Run", font_size=32, bold=True, color=WHITE)

# Step boxes
steps = [
    ("1", "Install Ollama", "ollama pull qwen3.5:4b\nollama pull qwen3.5:9b"),
    ("2", "Install Python", "pip install -r requirements.txt\n(3~10분 소요)"),
    ("3", "Launch App", "streamlit run app.py\n자동으로 브라우저 열림"),
    ("4", "Upload PDF", "Paper Lab 탭 → PDF 업로드\n→ 자동 인덱싱"),
    ("5", "Ask Questions", "채팅 입력 → Agent Pipeline\n→ 답변 + 출처 확인"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.8 + i * 2.4)
    y = Inches(1.6)
    # Number circle
    add_circle(slide, x + Inches(0.8), y, Inches(0.5), fill_color=ACCENT)
    add_textbox(slide, x + Inches(0.8), y + Inches(0.05), Inches(0.5), Inches(0.4),
                num, font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, x, y + Inches(0.7), Inches(2.2), Inches(0.4),
                title, font_size=14, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    # Description
    desc_lines = desc.split("\n")
    formatted = [(l, False, 11, MED_GRAY) for l in desc_lines]
    add_multiline_textbox(slide, x + Inches(0.1), y + Inches(1.1), Inches(2.0), Inches(1.0),
                          formatted, alignment=PP_ALIGN.CENTER)
    # Arrow between steps
    if i < len(steps) - 1:
        add_arrow_right(slide, x + Inches(2.0), y + Inches(0.6), Inches(0.3), Inches(0.2), fill_color=MED_GRAY)

# Bottom: Sample query box
add_rounded_rect(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.3), fill_color=LIGHT_GRAY, line_color=ACCENT, line_width=2)
add_textbox(slide, Inches(1.2), Inches(4.6), Inches(4), Inches(0.4),
            "💡 Sample Queries", font_size=18, bold=True, color=ACCENT)

samples = [
    "•  \"What does this paper say about TCAD calibration?\"",
    "•  \"Summarize the key findings on GAA FET\"",
    "•  \"Search for TCAD machine learning papers on Arxiv\"",
    "•  \"Write a report comparing ML methods in my papers\"",
]
formatted = [(s, False, 14, DARK_GRAY) for s in samples]
add_multiline_textbox(slide, Inches(1.2), Inches(5.1), Inches(10.5), Inches(1.5), formatted)


# ═══════════════════════════════════════════════════════════
# SLIDE 11 — CONCLUSION
# ═══════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)

# Decorative bar
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), fill_color=ACCENT3)

# Title
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
            "Thank You", font_size=42, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
            "Research Agent — TCAD + Semiconductor Device Paper Assistant",
            font_size=22, color=ACCENT3, alignment=PP_ALIGN.CENTER)

# Divider
add_rect(slide, Inches(4.5), Inches(3.3), Inches(4), Inches(0.03), fill_color=ACCENT3)

# Key takeaways
takeaways = [
    "📤 PDF Ingestion → BGE-m3 Embedding → FAISS Index",
    "🤖 Multi-Agent Pipeline (Supervisor → Paper/Arxiv/Report)",
    "💬 RAG Q&A with Source Citation",
    "🔍 Arxiv Search with LLM Summarization",
    "📋 Comprehensive Report Generation",
    "📝 체계적인 Trial & Error 문서화 (Final Report 포함)",
]
formatted = [(t, False, 17, WHITE) for t in takeaways]
add_multiline_textbox(slide, Inches(2.5), Inches(3.6), Inches(8), Inches(2.8), formatted, alignment=PP_ALIGN.LEFT)

# Bottom info
add_textbox(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
            "Fast Campus Gen AI Intensive Course  |  2026. 06. 29 — 07. 10",
            font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ── Save ──

output_path = "C:\\Users\\User\\Documents\\GenAI\\research-agent\\Research_Agent_Final_Presentation.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
